"""
版式分析器 - 智能识别 PPT 模板版式
"""

import logging
from dataclasses import dataclass, field
from typing import List, Dict, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

logger = logging.getLogger(__name__)


@dataclass
class TextCapacity:
    """文本容量估计"""
    title_max_chars: int = 100
    body_max_lines: int = 10
    bullet_max_items: int = 6


@dataclass
class LayoutFeatures:
    """版式特征"""
    slide_index: int
    placeholder_types: List[str] = field(default_factory=list)
    text_shape_count: int = 0
    image_placeholder_count: int = 0
    column_count: int = 1
    has_title: bool = False
    has_body: bool = False
    title_position: str = "top"  # top, center, left, right
    layout_symmetry: str = "asymmetric"  # symmetric, asymmetric
    visual_weight: Dict[str, float] = field(default_factory=dict)
    estimated_capacity: Optional[TextCapacity] = None
    layout_type: str = "unknown"
    shape_positions: List[Dict] = field(default_factory=list)


@dataclass
class LayoutProfile:
    """版式档案"""
    slide_index: int
    features: LayoutFeatures
    layout_type: str
    confidence: float
    description: str


class LayoutAnalyzer:
    """智能识别 PPT 模板版式"""
    
    def __init__(self):
        """初始化版式分析器"""
        self.layout_profiles: List[LayoutProfile] = []
        self.presentation: Optional[Presentation] = None
    
    def analyze(self, prs: Presentation) -> List[LayoutProfile]:
        """
        分析所有幻灯片版式
        
        Args:
            prs: Presentation 对象
            
        Returns:
            版式档案列表
        """
        self.layout_profiles = []
        self.presentation = prs
        
        for idx, slide in enumerate(prs.slides):
            logger.debug(f"分析幻灯片 {idx}/{len(prs.slides)}")
            
            # 提取特征
            features = self._extract_features(slide, idx)
            
            # 分类版式
            layout_type = self._classify_layout(features)
            
            # 创建档案
            profile = self._create_profile(slide, features, layout_type)
            self.layout_profiles.append(profile)
        
        logger.info(f"版式分析完成：{len(self.layout_profiles)}个版式")
        return self.layout_profiles
    
    def _extract_features(self, slide, slide_index: int) -> LayoutFeatures:
        """
        提取版式特征
        
        Args:
            slide: 幻灯片对象
            slide_index: 幻灯片索引
            
        Returns:
            LayoutFeatures 对象
        """
        features = LayoutFeatures(slide_index=slide_index)
        
        shape_positions = []
        text_shapes = []
        
        for shape in slide.shapes:
            shape_info = {
                'type': shape.shape_type,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'center_x': shape.left + shape.width / 2,
                'center_y': shape.top + shape.height / 2
            }
            shape_positions.append(shape_info)
            
            # 检查占位符
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                placeholder_type = shape.placeholder_format.type
                features.placeholder_types.append(PP_PLACEHOLDER(placeholder_type).name)
                
                if placeholder_type in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                    features.has_title = True
                    features.title_position = self._detect_title_position(shape, slide)
                
                if placeholder_type == PP_PLACEHOLDER.BODY:
                    features.has_body = True
                
                if placeholder_type == PP_PLACEHOLDER.OBJECT:
                    features.image_placeholder_count += 1
            
            # 统计文本形状
            if hasattr(shape, "text") and shape.text.strip():
                features.text_shape_count += 1
                text_shapes.append(shape_info)
        
        features.shape_positions = shape_positions
        
        # 检测列数
        features.column_count = self._detect_column_count(text_shapes, slide)
        
        # 检测对称性
        features.layout_symmetry = self._detect_symmetry(text_shapes, slide)
        
        # 计算视觉权重
        features.visual_weight = self._calculate_visual_weight(text_shapes, slide)
        
        # 估计文本容量
        features.estimated_capacity = self._estimate_capacity(slide, features)
        
        return features
    
    def _detect_title_position(self, shape, slide) -> str:
        """检测标题位置"""
        slide_width = self.presentation.slide_width if self.presentation else slide.slide_width
        slide_height = self.presentation.slide_height if self.presentation else slide.slide_height
        
        center_x = shape.left + shape.width / 2
        center_y = shape.top + shape.height / 2
        
        # 垂直居中
        if abs(center_y - slide_height / 2) < slide_height * 0.2:
            return "center"
        
        # 水平位置
        if center_x < slide_width * 0.3:
            return "left"
        elif center_x > slide_width * 0.7:
            return "right"
        else:
            return "top"
    
    def _detect_column_count(self, text_shapes: List[Dict], slide) -> int:
        """
        通过分析形状的 X 坐标分布检测列数
        """
        if not text_shapes:
            return 1
        
        # 收集 X 坐标
        x_centers = [shape['center_x'] for shape in text_shapes]
        
        if len(x_centers) < 2:
            return 1
        
        # 简单聚类：按 X 坐标分组
        x_centers_sorted = sorted(x_centers)
        columns = []
        current_column = [x_centers_sorted[0]]
        
        # 获取幻灯片宽度
        slide_width = self.presentation.slide_width if self.presentation else 9144000
        
        for i in range(1, len(x_centers_sorted)):
            gap = x_centers_sorted[i] - x_centers_sorted[i-1]
            # 如果间距较大，认为是新列
            if gap > slide_width * 0.15:
                columns.append(current_column)
                current_column = [x_centers_sorted[i]]
            else:
                current_column.append(x_centers_sorted[i])
        
        columns.append(current_column)
        
        # 限制列数范围
        return min(max(len(columns), 1), 4)
    
    def _detect_symmetry(self, text_shapes: List[Dict], slide) -> str:
        """检测版式对称性"""
        if len(text_shapes) < 2:
            return "asymmetric"
        
        slide_width = self.presentation.slide_width if self.presentation else 9144000
        slide_center = slide_width / 2
        
        # 检查形状是否关于中心对称分布
        left_shapes = [s for s in text_shapes if s['center_x'] < slide_center]
        right_shapes = [s for s in text_shapes if s['center_x'] > slide_center]
        
        if abs(len(left_shapes) - len(right_shapes)) <= 1:
            return "symmetric"
        
        return "asymmetric"
    
    def _calculate_visual_weight(self, text_shapes: List[Dict], slide) -> Dict[str, float]:
        """计算视觉权重分布"""
        slide_width = self.presentation.slide_width if self.presentation else 9144000
        slide_height = self.presentation.slide_height if self.presentation else 6858000
        slide_area = slide_width * slide_height
        
        weights = {
            'top': 0.0,
            'bottom': 0.0,
            'left': 0.0,
            'right': 0.0,
            'center': 0.0
        }
        
        for shape in text_shapes:
            shape_area = shape['width'] * shape['height']
            shape_ratio = shape_area / slide_area if slide_area > 0 else 0
            
            # 垂直分布
            shape_center_y = shape['center_y']
            
            if shape_center_y < slide_height * 0.33:
                weights['top'] += shape_ratio
            elif shape_center_y > slide_height * 0.66:
                weights['bottom'] += shape_ratio
            else:
                weights['center'] += shape_ratio
            
            # 水平分布
            shape_center_x = shape['center_x']
            
            if shape_center_x < slide_width * 0.33:
                weights['left'] += shape_ratio
            elif shape_center_x > slide_width * 0.66:
                weights['right'] += shape_ratio
        
        return weights
    
    def _estimate_capacity(self, slide, features: LayoutFeatures) -> TextCapacity:
        """估计文本容量"""
        capacity = TextCapacity()
        
        # 基于形状尺寸估计
        for shape_info in features.shape_positions:
            # 标题容量
            if 'TITLE' in str(shape_info.get('type', '')):
                # 根据宽度估计字符数（假设平均字符宽度约 20 EMU）
                capacity.title_max_chars = int(shape_info['width'] / 200000)  # 粗略估计
            
            # 正文容量
            if 'BODY' in str(shape_info.get('type', '')):
                # 根据高度估计行数（假设平均行高约 30 EMU）
                capacity.body_max_lines = int(shape_info['height'] / 300000)
        
        return capacity
    
    def _classify_layout(self, features: LayoutFeatures) -> str:
        """
        基于特征分类版式类型
        """
        # 规则引擎分类
        
        # 1. 封面页
        if (features.has_title and features.text_shape_count == 1 
            and features.title_position == 'center'):
            return 'TITLE_SLIDE'
        
        # 2. 章节页
        if (features.has_title and features.text_shape_count == 1 
            and features.title_position in ['left', 'top']):
            return 'SECTION_HEADER'
        
        # 3. 引用页
        if (features.text_shape_count == 2 
            and features.layout_symmetry == 'symmetric'
            and not features.has_body):
            return 'QUOTE'
        
        # 4. 三列对比
        if features.has_title and features.column_count == 3:
            return 'THREE_COLUMN'
        
        # 5. 双列对比
        if features.has_title and features.column_count == 2:
            return 'TWO_COLUMN'
        
        # 6. 左图右文
        if (features.has_title and features.image_placeholder_count >= 1 
            and features.visual_weight.get('right', 0) > 0.5):
            return 'IMAGE_RIGHT'
        
        # 7. 右图左文
        if (features.has_title and features.image_placeholder_count >= 1 
            and features.visual_weight.get('left', 0) > 0.5):
            return 'IMAGE_LEFT'
        
        # 8. 项目列表
        if features.has_title and features.has_body and features.column_count == 1:
            return 'BULLET_LIST'
        
        # 9. 上图下文
        if (features.has_title and features.visual_weight.get('top', 0) > 0.6):
            return 'IMAGE_TOP'
        
        # 10. 空白页
        if features.text_shape_count == 0:
            return 'BLANK'
        
        # 默认：带标题的内容页
        if features.has_title:
            return 'CONTENT_WITH_CAPTION'
        
        return 'CONTENT'
    
    def _create_profile(self, slide, features: LayoutFeatures, layout_type: str) -> LayoutProfile:
        """创建版式档案"""
        # 计算置信度
        confidence = self._calculate_confidence(features, layout_type)
        
        # 生成描述
        description = self._generate_description(features, layout_type)
        
        return LayoutProfile(
            slide_index=features.slide_index,
            features=features,
            layout_type=layout_type,
            confidence=confidence,
            description=description
        )
    
    def _calculate_confidence(self, features: LayoutFeatures, layout_type: str) -> float:
        """计算分类置信度"""
        # 基于特征匹配度计算
        confidence = 0.5  # 基础置信度
        
        # 特征越多，置信度越高
        if features.has_title:
            confidence += 0.1
        if features.has_body:
            confidence += 0.1
        if features.column_count > 1:
            confidence += 0.1
        
        # 特殊版式置信度更高
        if layout_type in ['TITLE_SLIDE', 'QUOTE', 'BLANK']:
            confidence += 0.2
        
        return min(confidence, 1.0)
    
    def _generate_description(self, features: LayoutFeatures, layout_type: str) -> str:
        """生成版式描述"""
        desc_parts = []
        
        if features.has_title:
            desc_parts.append(f"标题 ({features.title_position})")
        
        if features.has_body:
            desc_parts.append("正文")
        
        if features.column_count > 1:
            desc_parts.append(f"{features.column_count}列")
        
        if features.image_placeholder_count > 0:
            desc_parts.append(f"{features.image_placeholder_count}个图片位")
        
        desc = f"{layout_type}: {', '.join(desc_parts)}" if desc_parts else layout_type
        
        return desc
    
    def get_layouts_by_type(self, layout_type: str) -> List[LayoutProfile]:
        """
        根据类型获取版式
        
        Args:
            layout_type: 版式类型
            
        Returns:
            匹配的版式档案列表
        """
        return [p for p in self.layout_profiles if p.layout_type == layout_type]
    
    def get_best_layout(self, layout_type: str) -> Optional[LayoutProfile]:
        """
        获取某类型的最佳版式
        
        Args:
            layout_type: 版式类型
            
        Returns:
            最佳版式档案
        """
        layouts = self.get_layouts_by_type(layout_type)
        if not layouts:
            return None
        return max(layouts, key=lambda p: p.confidence)
    
    def export_profiles(self) -> List[Dict]:
        """
        导出版式档案为字典列表
        
        Returns:
            档案字典列表
        """
        return [
            {
                'slide_index': p.slide_index,
                'layout_type': p.layout_type,
                'confidence': p.confidence,
                'description': p.description,
                'features': {
                    'text_shape_count': p.features.text_shape_count,
                    'column_count': p.features.column_count,
                    'has_title': p.features.has_title,
                    'has_body': p.features.has_body,
                    'title_position': p.features.title_position,
                }
            }
            for p in self.layout_profiles
        ]
