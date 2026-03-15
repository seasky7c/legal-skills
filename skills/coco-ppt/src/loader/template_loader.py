"""
模板加载器 - 加载和验证 PPTX 模板
"""

import logging
from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from dataclasses import dataclass

logger = logging.getLogger(__name__)


@dataclass
class TemplateValidationResult:
    """模板验证结果"""
    is_valid: bool
    slide_count: int
    layout_count: int
    errors: List[str]
    warnings: List[str]


class TemplateLoader:
    """加载和验证 PPTX 模板"""
    
    def __init__(self, template_path: str):
        """
        初始化模板加载器
        
        Args:
            template_path: PPTX 模板文件路径
        """
        self.template_path = Path(template_path)
        self.presentation: Optional[Presentation] = None
        self.validation_result: Optional[TemplateValidationResult] = None
    
    def load(self) -> Presentation:
        """
        加载 PPTX 模板文件
        
        Returns:
            Presentation 对象
            
        Raises:
            FileNotFoundError: 模板文件不存在
            ValueError: 文件格式不正确
        """
        if not self.template_path.exists():
            raise FileNotFoundError(f"模板文件不存在：{self.template_path}")
        
        if self.template_path.suffix.lower() != '.pptx':
            raise ValueError(f"文件必须是 PPTX 格式：{self.template_path}")
        
        logger.info(f"加载模板：{self.template_path}")
        self.presentation = Presentation(self.template_path)
        
        # 自动验证
        self.validate()
        
        return self.presentation
    
    def validate(self) -> TemplateValidationResult:
        """
        验证模板是否有效
        
        Returns:
            TemplateValidationResult 验证结果
        """
        if self.presentation is None:
            self.load()
        
        errors = []
        warnings = []
        
        # 检查幻灯片数量
        slide_count = len(self.presentation.slides)
        if slide_count == 0:
            errors.append("模板中没有幻灯片")
        elif slide_count < 2:
            warnings.append(f"模板只有{slide_count}张幻灯片，可能不够用")
        
        # 检查版式数量
        layout_count = len(self.presentation.slide_layouts)
        if layout_count == 0:
            errors.append("模板中没有版式")
        elif layout_count < 3:
            warnings.append(f"模板只有{layout_count}个版式，可能缺乏多样性")
        
        # 检查是否有标题版式
        has_title_layout = any(
            layout.name and 'title' in layout.name.lower()
            for layout in self.presentation.slide_layouts
        )
        if not has_title_layout:
            warnings.append("模板可能没有标题版式")
        
        self.validation_result = TemplateValidationResult(
            is_valid=len(errors) == 0,
            slide_count=slide_count,
            layout_count=layout_count,
            errors=errors,
            warnings=warnings
        )
        
        logger.info(f"模板验证完成：{slide_count}张幻灯片，{layout_count}个版式")
        
        if errors:
            for error in errors:
                logger.error(f"验证错误：{error}")
        
        if warnings:
            for warning in warnings:
                logger.warning(f"验证警告：{warning}")
        
        return self.validation_result
    
    def get_slide_layouts(self) -> List[dict]:
        """
        提取所有幻灯片版式信息
        
        Returns:
            版式信息列表
        """
        if self.presentation is None:
            self.load()
        
        layouts = []
        for idx, layout in enumerate(self.presentation.slide_layouts):
            layouts.append({
                'index': idx,
                'name': layout.name,
                'master_slide_name': layout.master_slide.name if layout.master_slide else 'Unknown'
            })
        
        logger.debug(f"提取到{len(layouts)}个版式")
        return layouts
    
    def get_slides(self) -> List[dict]:
        """
        获取所有幻灯片的基本信息
        
        Returns:
            幻灯片信息列表
        """
        if self.presentation is None:
            self.load()
        
        slides = []
        for idx, slide in enumerate(self.presentation.slides):
            shape_count = len(slide.shapes)
            slides.append({
                'index': idx,
                'shape_count': shape_count,
                'notes': slide.has_notes_slide if hasattr(slide, 'has_notes_slide') else False
            })
        
        logger.debug(f"提取到{len(slides)}张幻灯片")
        return slides
