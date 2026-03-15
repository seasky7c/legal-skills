"""
内容应用器 - 将内容填充到幻灯片
"""

import logging
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Pt
from ..loader.content_loader import SlideContent
from ..analyzer.matcher import LayoutMatch

logger = logging.getLogger(__name__)


class ContentApplier:
    """应用内容到幻灯片"""
    
    def __init__(self, presentation: Presentation):
        """
        初始化内容应用器
        
        Args:
            presentation: Presentation 对象
        """
        self.presentation = presentation
        self.issues: List[Dict] = []
    
    def apply_all(self, matches: List[LayoutMatch]) -> List[Dict]:
        """
        应用所有内容到幻灯片
        
        Args:
            matches: 匹配结果列表
            
        Returns:
            问题列表
        """
        self.issues = []
        
        logger.info(f"开始应用内容：{len(matches)}张幻灯片")
        
        for idx, match in enumerate(matches):
            if idx < len(self.presentation.slides):
                slide = self.presentation.slides[idx]
                content = match.requirement.content
                self._apply_to_slide(slide, content, match)
        
        logger.info(f"内容应用完成，发现{len(self.issues)}个问题")
        return self.issues
    
    def _apply_to_slide(self, slide: Slide, content: SlideContent, match: LayoutMatch):
        """
        应用内容到单张幻灯片
        
        Args:
            slide: 幻灯片对象
            content: 幻灯片内容
            match: 匹配结果
        """
        logger.debug(f"应用内容到幻灯片 {match.slide_index}: {content.title or '(无标题)'}")
        
        # 根据版式类型应用内容
        layout_type = match.matched_layout.layout_type
        
        if layout_type == 'TITLE_SLIDE':
            self._apply_title_slide(slide, content)
        elif layout_type == 'SECTION_HEADER':
            self._apply_section_header(slide, content)
        elif layout_type == 'BULLET_LIST':
            self._apply_bullet_list(slide, content)
        elif layout_type == 'QUOTE':
            self._apply_quote(slide, content)
        elif layout_type in ['TWO_COLUMN', 'THREE_COLUMN']:
            self._apply_multi_column(slide, content, layout_type)
        else:
            self._apply_generic(slide, content)
    
    def _apply_title_slide(self, slide: Slide, content: SlideContent):
        """应用封面页"""
        shapes = self._get_shapes_by_type(slide, [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.SUBTITLE, PP_PLACEHOLDER.CENTER_TITLE])
        
        for shape in shapes:
            ptype = shape.placeholder_format.type if hasattr(shape, 'placeholder_format') else None
            
            if ptype in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                self._fill_shape(shape, content.title)
            elif ptype == PP_PLACEHOLDER.SUBTITLE:
                self._fill_shape(shape, content.subtitle)
    
    def _apply_section_header(self, slide: Slide, content: SlideContent):
        """应用章节页"""
        shapes = self._get_shapes_by_type(slide, [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE])
        
        for shape in shapes:
            self._fill_shape(shape, content.title)
    
    def _apply_bullet_list(self, slide: Slide, content: SlideContent):
        """应用项目列表"""
        shapes = self._get_shapes_by_type(slide, [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.BODY])
        
        for shape in shapes:
            ptype = shape.placeholder_format.type if hasattr(shape, 'placeholder_format') else None
            
            if ptype in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                self._fill_shape(shape, content.title)
            elif ptype == PP_PLACEHOLDER.BODY:
                if content.items:
                    self._fill_with_bullets(shape, content.items)
                else:
                    self._fill_shape(shape, content.content)
    
    def _apply_quote(self, slide: Slide, content: SlideContent):
        """应用引用页"""
        shapes = slide.shapes
        
        # 查找文本形状
        text_shapes = [s for s in shapes if hasattr(s, 'text') and s.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER]
        
        if len(text_shapes) >= 1:
            self._fill_shape(text_shapes[0], content.quote)
        
        if len(text_shapes) >= 2 and content.author:
            self._fill_shape(text_shapes[1], f"— {content.author}")
    
    def _apply_multi_column(self, slide: Slide, content: SlideContent, layout_type: str):
        """应用多列布局"""
        shapes = self._get_shapes_by_type(slide, [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT])
        
        title_applied = False
        for shape in shapes:
            ptype = shape.placeholder_format.type if hasattr(shape, 'placeholder_format') else None
            
            if ptype in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE] and not title_applied:
                self._fill_shape(shape, content.title)
                title_applied = True
            elif ptype == PP_PLACEHOLDER.BODY:
                if content.left_content and not hasattr(self, '_left_filled'):
                    self._fill_shape(shape, content.left_content.get('content', ''))
                    self._left_filled = True
                elif content.right_content:
                    self._fill_shape(shape, content.right_content.get('content', ''))
                elif content.content:
                    self._fill_shape(shape, content.content)
    
    def _apply_generic(self, slide: Slide, content: SlideContent):
        """应用通用内容"""
        shapes = self._get_shapes_by_type(slide, [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.CENTER_TITLE])
        
        title_applied = False
        for shape in shapes:
            ptype = shape.placeholder_format.type if hasattr(shape, 'placeholder_format') else None
            
            if ptype in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE] and not title_applied:
                self._fill_shape(shape, content.title)
                title_applied = True
            elif ptype == PP_PLACEHOLDER.BODY:
                if content.items:
                    self._fill_with_bullets(shape, content.items)
                else:
                    self._fill_shape(shape, content.content)
            elif not title_applied and content.title:
                self._fill_shape(shape, content.title)
                title_applied = True
    
    def _get_shapes_by_type(self, slide: Slide, placeholder_types: List) -> List[BaseShape]:
        """
        根据占位符类型获取形状
        
        Args:
            slide: 幻灯片
            placeholder_types: 占位符类型列表
            
        Returns:
            形状列表
        """
        shapes = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if shape.placeholder_format.type in placeholder_types:
                    shapes.append(shape)
        return shapes
    
    def _fill_shape(self, shape: BaseShape, text: str):
        """
        填充形状内容
        
        Args:
            shape: 形状对象
            text: 文本内容
        """
        if not text or not hasattr(shape, 'text'):
            return
        
        try:
            # 清除现有文本
            if shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = ""
                
                # 设置新文本
                if shape.text_frame.paragraphs:
                    shape.text_frame.paragraphs[0].text = text
                else:
                    shape.text = text
            
            logger.debug(f"填充形状：{text[:50]}...")
        except Exception as e:
            logger.warning(f"填充形状失败：{e}")
            self.issues.append({
                'type': 'fill_error',
                'message': str(e),
                'text': text[:100]
            })
    
    def _fill_with_bullets(self, shape: BaseShape, items: List[str]):
        """
        填充项目列表
        
        Args:
            shape: 形状对象
            items: 项目列表
        """
        if not items or not hasattr(shape, 'text_frame'):
            return
        
        try:
            tf = shape.text_frame
            tf.clear()
            
            for i, item in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = item
                p.level = 0
                
                # 设置项目符号
                from pptx.enum.text import PP_ALIGN
                from pptx.util import Pt
                
                # 注意：python-pptx 的项目符号支持有限
                # 这里简单设置文本，实际项目符号可能需要手动调整
            
            logger.debug(f"填充项目列表：{len(items)}项")
        except Exception as e:
            logger.warning(f"填充项目列表失败：{e}")
            self.issues.append({
                'type': 'bullet_error',
                'message': str(e),
                'items_count': len(items)
            })
    
    def get_issues(self) -> List[Dict]:
        """获取所有问题"""
        return self.issues
