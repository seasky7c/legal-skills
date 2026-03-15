"""
验证器 - 验证生成质量
"""

import logging
from typing import List, Dict
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


class ValidationResult:
    """验证结果"""
    
    def __init__(self):
        self.is_valid = True
        self.errors: List[Dict] = []
        self.warnings: List[Dict] = []
        self.info: List[Dict] = []
    
    def add_error(self, message: str, slide_index: int = -1, **kwargs):
        """添加错误"""
        self.is_valid = False
        self.errors.append({
            'type': 'error',
            'message': message,
            'slide_index': slide_index,
            **kwargs
        })
    
    def add_warning(self, message: str, slide_index: int = -1, **kwargs):
        """添加警告"""
        self.warnings.append({
            'type': 'warning',
            'message': message,
            'slide_index': slide_index,
            **kwargs
        })
    
    def add_info(self, message: str, slide_index: int = -1, **kwargs):
        """添加信息"""
        self.info.append({
            'type': 'info',
            'message': message,
            'slide_index': slide_index,
            **kwargs
        })
    
    def to_dict(self) -> Dict:
        """转换为字典"""
        return {
            'is_valid': self.is_valid,
            'error_count': len(self.errors),
            'warning_count': len(self.warnings),
            'info_count': len(self.info),
            'errors': self.errors,
            'warnings': self.warnings,
            'info': self.info
        }


class SlideValidator:
    """验证生成质量"""
    
    def __init__(self, presentation: Presentation):
        """
        初始化验证器
        
        Args:
            presentation: Presentation 对象
        """
        self.presentation = presentation
    
    def validate(self) -> ValidationResult:
        """
        验证整个演示文稿
        
        Returns:
            ValidationResult 对象
        """
        result = ValidationResult()
        
        logger.info(f"开始验证：{len(self.presentation.slides)}张幻灯片")
        
        # 基础验证
        if len(self.presentation.slides) == 0:
            result.add_error("演示文稿没有幻灯片")
            return result
        
        # 逐张验证
        for idx, slide in enumerate(self.presentation.slides):
            slide_result = self._validate_slide(slide, idx)
            result.errors.extend(slide_result.errors)
            result.warnings.extend(slide_result.warnings)
            result.info.extend(slide_result.info)
        
        if result.errors:
            result.is_valid = False
        
        logger.info(f"验证完成：{len(result.errors)}错误，{len(result.warnings)}警告")
        return result
    
    def _validate_slide(self, slide: Slide, index: int) -> ValidationResult:
        """
        验证单张幻灯片
        
        Args:
            slide: 幻灯片
            index: 索引
            
        Returns:
            ValidationResult 对象
        """
        result = ValidationResult()
        
        # 检查空形状
        empty_shapes = self._check_empty_shapes(slide)
        for shape in empty_shapes:
            result.add_warning(
                f"形状为空",
                slide_index=index,
                shape_id=id(shape)
            )
        
        # 检查文本溢出 (简化版)
        overflow_shapes = self._check_overflow(slide)
        for shape_info in overflow_shapes:
            result.add_warning(
                f"文本可能溢出",
                slide_index=index,
                shape_info=shape_info
            )
        
        # 检查标题
        has_title = self._check_has_title(slide)
        if not has_title and index > 0:  # 第一张可以是封面
            result.add_info(
                "幻灯片没有标题",
                slide_index=index
            )
        
        return result
    
    def _check_empty_shapes(self, slide: Slide) -> List[BaseShape]:
        """
        检查空形状
        
        Args:
            slide: 幻灯片
            
        Returns:
            空形状列表
        """
        empty = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if hasattr(shape, 'text') and not shape.text.strip():
                    empty.append(shape)
        return empty
    
    def _check_overflow(self, slide: Slide) -> List[Dict]:
        """
        检查文本溢出
        
        Args:
            slide: 幻灯片
            
        Returns:
            溢出形状信息列表
        """
        overflow = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                tf = shape.text_frame
                
                # 检查文本是否超出边界
                if tf.auto_size is False:  # 固定大小
                    text_height = self._estimate_text_height(tf)
                    if text_height > shape.height:
                        overflow.append({
                            'shape_id': id(shape),
                            'text_length': len(shape.text),
                            'estimated_height': text_height,
                            'actual_height': shape.height
                        })
        
        return overflow
    
    def _estimate_text_height(self, text_frame) -> int:
        """
        估算文本高度
        
        Args:
            text_frame: 文本框
            
        Returns:
            估算高度 (EMU)
        """
        # 简化估算：每行约 30 EMU
        total_lines = sum(len(p.runs) for p in text_frame.paragraphs if p.runs)
        return total_lines * 300000  # 粗略估算
    
    def _check_has_title(self, slide: Slide) -> bool:
        """
        检查是否有标题
        
        Args:
            slide: 幻灯片
            
        Returns:
            是否有标题
        """
        from pptx.enum.shapes import PP_PLACEHOLDER
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if hasattr(shape, 'placeholder_format'):
                    ptype = shape.placeholder_format.type
                    if ptype in [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]:
                        if hasattr(shape, 'text') and shape.text.strip():
                            return True
        
        return False
    
    def generate_report(self, result: ValidationResult) -> str:
        """
        生成验证报告
        
        Args:
            result: 验证结果
            
        Returns:
            报告文本
        """
        lines = [
            "=" * 60,
            "COCO-PPT 质量验证报告",
            "=" * 60,
            "",
            f"总幻灯片数：{len(self.presentation.slides)}",
            f"验证状态：{'✅ 通过' if result.is_valid else '❌ 失败'}",
            f"错误数：{len(result.errors)}",
            f"警告数：{len(result.warnings)}",
            f"信息数：{len(result.info)}",
            "",
        ]
        
        if result.errors:
            lines.append("-" * 60)
            lines.append("❌ 错误")
            lines.append("-" * 60)
            for error in result.errors:
                slide_info = f"(幻灯片 {error['slide_index'] + 1})" if error['slide_index'] >= 0 else ""
                lines.append(f"  {slide_info} {error['message']}")
            lines.append("")
        
        if result.warnings:
            lines.append("-" * 60)
            lines.append("⚠️  警告")
            lines.append("-" * 60)
            for warning in result.warnings:
                slide_info = f"(幻灯片 {warning['slide_index'] + 1})" if warning['slide_index'] >= 0 else ""
                lines.append(f"  {slide_info} {warning['message']}")
            lines.append("")
        
        if result.info:
            lines.append("-" * 60)
            lines.append("ℹ️  信息")
            lines.append("-" * 60)
            for info in result.info:
                slide_info = f"(幻灯片 {info['slide_index'] + 1})" if info['slide_index'] >= 0 else ""
                lines.append(f"  {slide_info} {info['message']}")
            lines.append("")
        
        lines.append("=" * 60)
        
        return "\n".join(lines)
